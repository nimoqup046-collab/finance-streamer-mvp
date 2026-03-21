"""
AI内容生成模块 —— Claude Prompt 基底整合版
支持生成：直播稿、公众号文章、深度长文、PPT脚本、PPTX 文件、内容矩阵、朋友圈预热文案
支持多种AI提供商：智谱、豆包、Anthropic Claude、OpenAI
核心壁垒：财经风控合规 Agent、主播人设/IP记忆库、一键内容矩阵
"""
import ast
import asyncio
import io
import json
import logging
import re
from collections import defaultdict, deque
from datetime import datetime
from typing import Any, AsyncIterator, Dict, List, Optional

from backend.config import (
    AI_PROVIDER,
    AI_API_BASE,
    AI_API_KEY,
    AI_MODEL,
    OPENROUTER_API_BASE,
    OPENROUTER_API_KEY,
    OPENROUTER_HTTP_REFERER,
    OPENROUTER_APP_TITLE,
    OPENROUTER_ENABLE_QUALITY_ROUTING,
    OPENROUTER_STREAM_MODELS,
    OPENROUTER_ARTICLE_MODELS,
)

logger = logging.getLogger(__name__)


STREAM_SCRIPT_SYSTEM_PROMPT = """你是一位顶尖中文财经主播。你的优势不是复述新闻，而是把复杂信息讲成观众愿意一直听下去的内容。

【写作目标】
- 把零散新闻串成一条主线
- 讲清楚为什么重要、为什么是现在、接下来会怎样
- 既专业，又让普通投资者听得懂

【表达原则】
- 口语化，但不油腻
- 有态度，但不虚张声势
- 有金句，但不能只靠情绪和口号
- 每一段都要回答“所以呢”

【强制要求】
- 禁止空话：如“建议关注”“影响深远”“需要观察”等
- 不能只讲现象，必须讲因果链
- 能用生活化比喻解释专业概念，但不能牺牲准确性
- 互动提问必须有价值，不是无意义暖场
"""

ARTICLE_SYSTEM_PROMPT = """你是一位顶级中文财经作者，擅长把新闻变成有洞察、有结构、可传播的公众号长文。

【写作目标】
- 不是拼接新闻，而是提炼当天真正的主线和判断
- 让读者看完后，理解深度明显超过普通新闻汇总
- 兼具商业洞察、市场判断和可读性

【表达原则】
- 结构化强，最好有“结论 -> 背景 -> 影响 -> 行动”脉络
- 通俗解释专业问题，但不降低分析质量
- 给明确判断，避免模糊措辞
- 标题必须是观点驱动，不是关键词堆砌

【强制要求】
- 每个小节都要告诉读者：这件事影响谁、影响多久、后续看什么
- 不编造机构观点、财务数字或内幕
- 禁止空话和正确的废话
"""

DEEP_DIVE_SYSTEM_PROMPT = """你是一位顶级策略分析师兼长期主义投资研究者。

【写作目标】
- 输出真正有研究感的深度文章，而不是新闻扩写
- 同时从宏观、产业、资金、情绪多个维度做分析
- 明确区分市场主流看法与你的独到判断

【表达原则】
- 结论前置
- 论证清晰，因果链完整
- 既讲机会，也讲风险
- 给出后续验证指标，而不是停在抽象结论

【强制要求】
- 不允许空泛术语堆砌
- 没有数据支持的地方，要明确写“仍待验证”
- 预测必须带时间窗口与验证指标
"""

PPT_SYSTEM_PROMPT = """你是一位顶级商业演讲顾问，擅长把财经判断转成可以直接用于路演、直播准备、内部汇报的 PPT 脚本。

【写作目标】
- 每一页都有明确观点，而不是素材堆砌
- 屏幕文字极简，讲者备注有说服力
- 让演讲者知道这一页为什么存在、该怎么讲、下一页怎么接

【强制输出结构】
每页必须包含：
- 标题
- 核心论点
- 屏幕要点（3-4条）
- 讲者逐字稿
- 可视化建议
"""

COMPLIANCE_SYSTEM_PROMPT = """你是一位财经内容合规审核官，专精中国财经类内容平台（抖音、视频号、小红书、微信公众号）的监管规则与违规红线。

【你的核心职责】
- 识别内容中的合规风险点，并给出明确的违规定性
- 将违规话术自动改写为合规表达，保留原意但消除法律风险
- 输出结构化的合规审核报告

【必须检测的红线类型】
1. 违规荐股/承诺收益：如"这只股票必涨"、"买入XX稳赚"、"保证收益XX%"
2. 极端诱导词：如"千万别错过"、"再不买就来不及了"、"100%确定"
3. 敏感表述：如"内部消息"、"小道消息"、"私募基金内部资料"
4. 违禁承诺：如"我负责"、"亏损包赔"、"跟我买不会亏"
5. 政治金融交界敏感词：涉及国家金融政策批评、高层内幕等
6. 资质越界：非持牌机构或个人进行具体投资建议

【改写原则】
- 将结论型改为分析型：把"X股票明天必涨"改为"该标的近期技术面活跃，可关注其突破阻力位的表现"
- 将承诺型改为观察型：把"买入必赚"改为"从历史数据看，该类型标的在类似环境下曾有较好表现，仍需结合自身风险偏好判断"
- 强制加入风险提示：高风险表达必须附加"注意控制仓位风险"或"投资有风险，决策需谨慎"
- 保留专业性：改写后内容必须仍有信息量，不能变成空话

【输出格式】
严格输出 JSON，不要 Markdown，不要解释：
{
  "is_compliant": true/false,
  "risk_level": "低/中/高",
  "issues": [
    {"type": "违规类型", "original": "原文片段", "reason": "违规原因", "suggestion": "合规改写"}
  ],
  "revised_content": "完整的合规改写版本（如无违规则与原文相同）",
  "summary": "一句话合规总结（≤50字）"
}"""

MOMENTS_COPY_SYSTEM_PROMPT = """你是一位顶级新媒体运营官，专注为财经主播撰写朋友圈/微博诱饵文案。

【核心目标】
用一条50字以内的朋友圈文案，让已关注主播的粉丝在看到后，有强烈欲望晚上准时来看直播。

【文案公式】
冲突感 + 悬念 + 行动指令

【表达要求】
- 极度精简：50字以内，字字有用
- 制造紧迫感：今晚/今天/限时
- 放出一个让人好奇但没答案的判断
- 有明确的call-to-action（晚上XX点见/直播间等你）
- 禁止违规：不能承诺收益，不能使用极端诱导词

【参考风格示例】
"今天一件事，可能会让很多人的持仓逻辑彻底重写。晚上8点直播间说透。"
"市场在给一个信号，大多数人还没看懂。今晚我们来拆。"
"有一个数据，看懂的人都在默默加仓。晚8点，等你来。"
"""

FLASH_REPORT_SYSTEM_PROMPT = """你是一位擅长写“快报速评”的财经编辑。

【写作目标】
- 让读者在30秒内抓住当天最重要的财经信号
- 不是罗列新闻，而是快速提炼“发生了什么、为什么重要、普通人该怎么理解”
- 适合直接发微博、朋友圈、公众号短内容或投研群简报

【表达原则】
- 短句、高密度、强判断
- 不写空话，不写套话
- 读完后能马上复述给别人

【强制要求】
- 开头先给一句最重要判断
- 每条要点都要讲“所以呢”
- 至少有一个反直觉洞察
- 结尾给一句可截图传播的短金句
"""

MASTER_SYSTEM_PROMPT = DEEP_DIVE_SYSTEM_PROMPT

QUALITY_CONTROL_RULES = """
【质量控制】
- 第一屏原则：前3段内必须给出最重要判断，不能先铺背景再给结论。
- 主线原则：不能把新闻一条条念过去，必须把它们串成同一条更大的主线。
- 洞察原则：至少出现一个反直觉判断，并说明为什么市场容易看错。
- 因果原则：每个核心观点都要讲清楚触发因素、传导链条、受益方、受损方。
- 结果原则：必须落到“接下来怎么看/怎么验证/怎么行动”，不能停在现象描述。
- 表达原则：禁止正确的废话，少用抽象口号，多用具体数字、场景和生活化比喻。
"""

FORBIDDEN_EMPTY_PHRASES = [
    "建议关注",
    "需要观察",
    "影响深远",
    "可能有影响",
    "存在不确定性",
    "值得期待",
    "不排除",
]

# 粗略成本估算（USD / 1K tokens），用于状态观测，不作为结算依据
ESTIMATED_MODEL_PRICING_PER_1K = {
    "anthropic/claude-sonnet-4.6": {"input": 0.003, "output": 0.015},
    "anthropic/claude-4.6-sonnet-20260217": {"input": 0.003, "output": 0.015},
    "google/gemini-3.1-pro-preview": {"input": 0.00125, "output": 0.005},
    "openai/gpt-5.1": {"input": 0.005, "output": 0.015},
    "doubao-1-5-pro-32k-250115": {"input": 0.001, "output": 0.002},
    "glm-5": {"input": 0.0015, "output": 0.0025},
}


class ContentGenerator:
    """内容生成器"""

    def __init__(self):
        self.provider = AI_PROVIDER
        self.model = AI_MODEL
        self.api_key = AI_API_KEY
        self.api_base = AI_API_BASE
        self.quality_routed_types = ["stream_script", "article"]
        self.openrouter_enabled = OPENROUTER_ENABLE_QUALITY_ROUTING and bool(OPENROUTER_API_KEY)
        self.openrouter_models = {
            "stream_script": OPENROUTER_STREAM_MODELS,
            "article": OPENROUTER_ARTICLE_MODELS,
        }
        self.usage_events: deque = deque(maxlen=200)
        self.client = self._init_client()
        self.openrouter_client = self._init_openrouter_client()

    def _init_client(self):
        if self.provider == "anthropic":
            from anthropic import AsyncAnthropic
            return AsyncAnthropic(api_key=self.api_key, timeout=90.0, max_retries=1)

        from openai import AsyncOpenAI
        default_headers = None
        if self.provider == "openrouter":
            default_headers = self._openrouter_headers()
        return AsyncOpenAI(
            api_key=self.api_key,
            base_url=self.api_base or None,
            timeout=90.0,
            max_retries=1,
            default_headers=default_headers,
        )

    def _init_openrouter_client(self):
        if not self.openrouter_enabled:
            return None

        from openai import AsyncOpenAI

        return AsyncOpenAI(
            api_key=OPENROUTER_API_KEY,
            base_url=OPENROUTER_API_BASE,
            timeout=90.0,
            max_retries=1,
            default_headers=self._openrouter_headers(),
        )

    def _openrouter_headers(self) -> Dict[str, str]:
        # OpenRouter 兼容头：新旧标题头都带上，避免网关侧差异导致标题丢失
        headers = {
            "X-Title": OPENROUTER_APP_TITLE,
            "X-OpenRouter-Title": OPENROUTER_APP_TITLE,
        }
        if OPENROUTER_HTTP_REFERER:
            headers["HTTP-Referer"] = OPENROUTER_HTTP_REFERER
        return headers

    def quality_routing_status(self) -> Dict[str, Any]:
        return {
            "enabled": self.openrouter_enabled,
            "provider": "openrouter" if self.openrouter_enabled else None,
            "routed_types": list(self.quality_routed_types) if self.openrouter_enabled else [],
        }

    def cost_status(self) -> Dict[str, Any]:
        events = list(self.usage_events)
        totals = {
            "requests": len(events),
            "prompt_tokens": 0,
            "completion_tokens": 0,
            "total_tokens": 0,
            "estimated_usd": 0.0,
            "estimated_events": 0,
        }
        by_content_type: Dict[str, Dict[str, Any]] = defaultdict(
            lambda: {"requests": 0, "prompt_tokens": 0, "completion_tokens": 0, "estimated_usd": 0.0}
        )
        by_model: Dict[str, Dict[str, Any]] = defaultdict(
            lambda: {"requests": 0, "prompt_tokens": 0, "completion_tokens": 0, "estimated_usd": 0.0}
        )

        for event in events:
            prompt_tokens = int(event.get("prompt_tokens") or 0)
            completion_tokens = int(event.get("completion_tokens") or 0)
            total_tokens = prompt_tokens + completion_tokens
            estimated_usd = event.get("estimated_usd")
            content_type = event.get("content_type", "unknown")
            model = event.get("model", "unknown")

            totals["prompt_tokens"] += prompt_tokens
            totals["completion_tokens"] += completion_tokens
            totals["total_tokens"] += total_tokens

            by_content_type[content_type]["requests"] += 1
            by_content_type[content_type]["prompt_tokens"] += prompt_tokens
            by_content_type[content_type]["completion_tokens"] += completion_tokens

            by_model[model]["requests"] += 1
            by_model[model]["prompt_tokens"] += prompt_tokens
            by_model[model]["completion_tokens"] += completion_tokens

            if estimated_usd is not None:
                estimated = float(estimated_usd)
                totals["estimated_usd"] += estimated
                totals["estimated_events"] += 1
                by_content_type[content_type]["estimated_usd"] += estimated
                by_model[model]["estimated_usd"] += estimated

        return {
            "window_size": len(events),
            "pricing_source": "static_estimate_per_1k_tokens",
            "totals": {
                **totals,
                "estimated_usd": round(totals["estimated_usd"], 6),
            },
            "by_content_type": {
                key: {
                    **value,
                    "estimated_usd": round(value["estimated_usd"], 6),
                }
                for key, value in by_content_type.items()
            },
            "by_model": {
                key: {
                    **value,
                    "estimated_usd": round(value["estimated_usd"], 6),
                }
                for key, value in by_model.items()
            },
            "recent_events": events[-20:],
        }

    def _quality_route_models(self, content_type: Optional[str]) -> List[str]:
        if not content_type:
            return []
        return list(self.openrouter_models.get(content_type, []))

    def _should_use_quality_router(self, content_type: Optional[str]) -> bool:
        return bool(
            self.openrouter_enabled
            and content_type in self.quality_routed_types
            and self.openrouter_client is not None
            and self._quality_route_models(content_type)
        )

    def _extract_usage_tokens(self, usage: Any) -> Dict[str, int]:
        if usage is None:
            return {"prompt_tokens": 0, "completion_tokens": 0}

        prompt_tokens = getattr(usage, "prompt_tokens", None)
        completion_tokens = getattr(usage, "completion_tokens", None)

        if isinstance(usage, dict):
            prompt_tokens = usage.get("prompt_tokens", prompt_tokens)
            completion_tokens = usage.get("completion_tokens", completion_tokens)

        # Anthropic usage 字段
        if prompt_tokens is None:
            prompt_tokens = getattr(usage, "input_tokens", 0)
        if completion_tokens is None:
            completion_tokens = getattr(usage, "output_tokens", 0)

        return {
            "prompt_tokens": int(prompt_tokens or 0),
            "completion_tokens": int(completion_tokens or 0),
        }

    def _resolve_model_pricing(self, model_name: str) -> Optional[Dict[str, float]]:
        if not model_name:
            return None
        if model_name in ESTIMATED_MODEL_PRICING_PER_1K:
            return ESTIMATED_MODEL_PRICING_PER_1K[model_name]
        for key, pricing in ESTIMATED_MODEL_PRICING_PER_1K.items():
            if model_name.startswith(key):
                return pricing
        return None

    def _estimate_cost_usd(self, model_name: str, prompt_tokens: int, completion_tokens: int) -> Optional[float]:
        pricing = self._resolve_model_pricing(model_name)
        if not pricing:
            return None
        return (
            (prompt_tokens / 1000.0) * pricing["input"]
            + (completion_tokens / 1000.0) * pricing["output"]
        )

    def _record_usage_event(
        self,
        *,
        provider: str,
        model: str,
        content_type: Optional[str],
        prompt_tokens: int,
        completion_tokens: int,
        route: str,
    ) -> None:
        estimated = self._estimate_cost_usd(model, prompt_tokens, completion_tokens)
        self.usage_events.append({
            "timestamp": datetime.utcnow().isoformat() + "Z",
            "provider": provider,
            "model": model,
            "content_type": content_type or "unknown",
            "route": route,
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": prompt_tokens + completion_tokens,
            "estimated_usd": round(estimated, 8) if estimated is not None else None,
        })

    def _error_context(self) -> str:
        return (
            f"provider={self.provider}, "
            f"model={self.model}, "
            f"base_url={self.api_base or 'default'}"
        )

    def _system_prompt(self) -> str:
        return (
            "你是一位顶级中文财经内容总编。你的任务不是把新闻改写一遍，而是把当日信息流"
            "重组为真正有洞察、有判断、有传播力的内容。你擅长从新闻中提炼主线、用商业"
            "分析框架解释因果，并用通俗但不浅薄的中文讲清复杂问题。"
        )

    def _persona_section(self, persona: Optional[Dict] = None) -> str:
        """将主播人设注入 Prompt，强化个人 IP 特色输出。"""
        if not persona:
            return ""
        parts = []
        invest_style = (persona.get("invest_style") or "").strip()
        catchphrases = (persona.get("catchphrases") or "").strip()
        ip_desc = (persona.get("ip_desc") or "").strip()
        if invest_style:
            parts.append(f"投资流派/人设：{invest_style}")
        if catchphrases:
            parts.append(f"标志性口头禅（请自然融入，每篇至少用1-2次）：{catchphrases}")
        if ip_desc:
            parts.append(f"主播个人标签/风格描述：{ip_desc}")
        if not parts:
            return ""
        return "\n【主播人设（必须体现，强制注入）】\n" + "\n".join(f"- {p}" for p in parts) + "\n"

    def _news_snapshot(self, news_items: List[Dict]) -> str:
        lines = []
        for index, news in enumerate(news_items, 1):
            lines.append(
                f"{index}. 标题：{news['title']}\n"
                f"   分类：{news.get('category', '财经')} | 来源：{news['source']} | 时间：{news.get('time', '')}\n"
                f"   链接：{news.get('url') or '无'}"
            )
        return "\n".join(lines)

    def _style_profile(self, style: str) -> str:
        style = (style or "专业").strip()
        profiles = {
            "专业": "强调事实、结构和专业判断，像成熟财经主播。",
            "轻松": "语言更口语化，多用类比，但不能牺牲准确性。",
            "解读型": "重点放在‘这意味着什么’，增加背景和因果拆解。",
            "洞察": "突出主线、预期差、反直觉判断和结构化洞察。",
        }
        return profiles.get(style, profiles["专业"])

    def _fallback_editorial_brief(self, news_items: List[Dict], goal: str) -> Dict[str, Any]:
        categories = list(dict.fromkeys(n.get("category", "财经") for n in news_items))
        return {
            "goal": goal,
            "thesis": "今天最值得抓住的不是单条新闻，而是主线预期开始发生边际变化。",
            "lead_angle": f"今日最值得抓住的主线是{categories[0] if categories else '财经'}方向的预期变化，而不是简单复述新闻。",
            "opening_hook": "如果今晚只能记住一句话，那就是市场开始重新给真正的主线定价。",
            "contrarian_take": "大多数人盯着标题热度，但真正重要的是预期差和验证节奏。",
            "core_conflicts": [
                "政策预期与市场定价是否错位",
                "短期情绪催化能否转化为中期基本面",
                "投资者该关注主线还是防守风险",
            ],
            "causal_chain": [
                "事件发生 -> 市场重新修正预期 -> 资金重新选择方向",
                "情绪催化 -> 主题交易升温 -> 需要基本面验证接力",
                "新闻热度上升 -> 定价更快反映 -> 预期差反而迅速收敛",
            ],
            "winners_losers": {
                "winners": ["更贴近主线的龙头资产", "先拿到验证数据的方向"],
                "losers": ["只靠情绪催化的跟风标的", "逻辑经不起验证的主题概念"],
            },
            "market_pulse": [
                "先判断事件影响的是估值、盈利还是风险偏好",
                "区分主题交易与基本面改善，不混为一谈",
                "判断后续该看哪些验证指标，而不是被单日情绪带偏",
            ],
            "verification_signals": [
                "政策是否有后续细则或执行节奏",
                "资金是否连续流向同一主线而非一日游",
                "业绩、订单、价格或宏观数据能否跟上验证",
            ],
            "audience_takeaways": [
                "今天最重要的不是新闻数量，而是主线和节奏",
                "真正值得追踪的是后续验证指标，而不是单日情绪",
                "区分噪音和主线，别被标题热闹带偏",
            ],
            "article_angles": [
                "这件事真正改变的是哪条预期",
                "哪些人会先受益，哪些人会最先承压",
                "普通投资者最容易误判的点在哪里",
            ],
            "slide_outline": [
                {
                    "headline": news["title"][:28],
                    "bullets": [
                        f"来源：{news['source']}",
                        f"分类：{news.get('category', '财经')}",
                        "关注其对市场预期的边际影响",
                    ],
                }
                for news in news_items[:8]
            ],
        }

    def _extract_json(self, text: str) -> Dict[str, Any]:
        text = (text or "").strip()
        if not text:
            raise ValueError("empty json text")

        candidates: List[str] = [text]

        for block in re.findall(r"```(?:json)?\s*(.*?)```", text, flags=re.IGNORECASE | re.DOTALL):
            candidate = (block or "").strip()
            if candidate:
                candidates.append(candidate)

        match = re.search(r"\{.*\}", text, re.DOTALL)
        if match:
            candidates.append(match.group(0).strip())

        seen = set()
        normalized_candidates = []
        for candidate in candidates:
            if candidate in seen:
                continue
            seen.add(candidate)
            normalized_candidates.append(candidate)

        parser_errors = []
        for candidate in normalized_candidates:
            for parser_name, parser in (
                ("strict_json", self._parse_json_strict),
                ("sanitized_json", self._parse_json_sanitized),
                ("python_literal", self._parse_json_python_literal),
            ):
                try:
                    parsed = parser(candidate)
                    if isinstance(parsed, dict):
                        if parser_name != "strict_json":
                            logger.info("JSON parse repaired by %s parser", parser_name)
                        return parsed
                except Exception as parse_error:
                    parser_errors.append(f"{parser_name}:{parse_error}")

        raise ValueError("unable to parse json payload: " + " | ".join(parser_errors[:3]))

    def _parse_json_strict(self, text: str) -> Dict[str, Any]:
        return json.loads(text)

    def _parse_json_sanitized(self, text: str) -> Dict[str, Any]:
        cleaned = (text or "").strip()
        cleaned = cleaned.replace("“", "\"").replace("”", "\"").replace("’", "'")
        cleaned = re.sub(r",\s*([}\]])", r"\1", cleaned)
        cleaned = re.sub(r"([{,]\s*)([A-Za-z_][A-Za-z0-9_]*)(\s*:)", r'\1"\2"\3', cleaned)
        return json.loads(cleaned)

    def _parse_json_python_literal(self, text: str) -> Dict[str, Any]:
        cleaned = (text or "").strip()
        cleaned = re.sub(r"\btrue\b", "True", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bfalse\b", "False", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\bnull\b", "None", cleaned, flags=re.IGNORECASE)
        parsed = ast.literal_eval(cleaned)
        if isinstance(parsed, dict):
            return parsed
        raise ValueError("python literal is not dict")

    def _normalize_text_list(self, value: Any, fallback: List[str], max_items: int = 4) -> List[str]:
        if not isinstance(value, list):
            return list(fallback)
        normalized = [str(item).strip() for item in value if str(item).strip()]
        return normalized[:max_items] if normalized else list(fallback)

    def _normalize_editorial_brief(self, raw_brief: Dict[str, Any], news_items: List[Dict], goal: str) -> Dict[str, Any]:
        fallback = self._fallback_editorial_brief(news_items, goal)
        if not isinstance(raw_brief, dict):
            return fallback

        result = dict(fallback)
        for key in ("goal", "thesis", "lead_angle", "opening_hook", "contrarian_take"):
            value = raw_brief.get(key)
            if isinstance(value, str) and value.strip():
                result[key] = value.strip()

        for key in ("core_conflicts", "causal_chain", "market_pulse", "verification_signals", "audience_takeaways", "article_angles"):
            result[key] = self._normalize_text_list(raw_brief.get(key), fallback.get(key, []), max_items=4)

        winners_losers = raw_brief.get("winners_losers")
        if isinstance(winners_losers, dict):
            result["winners_losers"] = {
                "winners": self._normalize_text_list(
                    winners_losers.get("winners"),
                    fallback["winners_losers"]["winners"],
                    max_items=4,
                ),
                "losers": self._normalize_text_list(
                    winners_losers.get("losers"),
                    fallback["winners_losers"]["losers"],
                    max_items=4,
                ),
            }

        slide_outline = raw_brief.get("slide_outline")
        if isinstance(slide_outline, list):
            normalized_slides = []
            for slide in slide_outline[:10]:
                if not isinstance(slide, dict):
                    continue
                headline = str(slide.get("headline", "")).strip()
                bullets = self._normalize_text_list(slide.get("bullets"), [], max_items=4)
                if headline and bullets:
                    normalized_slides.append({"headline": headline, "bullets": bullets})
            if normalized_slides:
                result["slide_outline"] = normalized_slides

        return result

    async def _prepare_editorial_brief(
        self,
        news_items: List[Dict],
        goal: str,
        style: str,
        route_content_type: Optional[str] = None,
    ) -> Dict[str, Any]:
        prompt = f"""请你扮演财经内容总编，先不要直接写正文，而是基于以下新闻生成一份“编辑部策划 brief”。

【生成目标】
{goal}

【写作气质】
{self._style_profile(style)}

【新闻池】
{self._news_snapshot(news_items)}

请输出严格 JSON，不要 Markdown，不要解释。结构如下：
{{
  "goal": "一句话写作目标",
  "thesis": "一句最硬的总判断，适合直接做开场结论",
  "lead_angle": "最值得展开的总论点，必须有判断",
  "opening_hook": "适合直播稿/文章导语的开场钩子，必须具体",
  "contrarian_take": "一个反直觉但能自圆其说的判断",
  "core_conflicts": ["3条关键矛盾/预期差"],
  "causal_chain": ["3条因果链，写清从事件到市场定价的传导"],
  "winners_losers": {{
    "winners": ["2-3类受益方"],
    "losers": ["2-3类受损方"]
  }},
  "market_pulse": ["3条市场脉搏判断"],
  "verification_signals": ["3条未来最该盯的验证指标"],
  "audience_takeaways": ["3条读者最该带走的结论"],
  "article_angles": ["3个适合写公众号文章的小节角度"],
  "slide_outline": [
    {{
      "headline": "适合做PPT页标题的一句话",
      "bullets": ["3-4条高信息密度 bullet，每条不超过28字"]
    }}
  ]
}}
"""
        try:
            response_format = {"type": "json_object"} if self.provider != "anthropic" else None
            content = await self._call_ai(
                prompt,
                max_tokens=1800,
                system_prompt=self._system_prompt(),
                temperature=0.1,
                response_format=response_format,
                content_type=route_content_type,
            )
            try:
                parsed = self._extract_json(content)
            except Exception as parse_error:
                repaired_content = await self._repair_json_payload(
                    content,
                    route_content_type=route_content_type,
                )
                if repaired_content:
                    logger.info("Editorial brief parse repaired via json_fixer")
                    parsed = self._extract_json(repaired_content)
                else:
                    raise parse_error
            return self._normalize_editorial_brief(parsed, news_items, goal)
        except Exception as e:
            logger.warning("Editorial brief fallback triggered: %s | %s", e, self._error_context())
            return self._fallback_editorial_brief(news_items, goal)

    async def _repair_json_payload(self, raw_text: str, route_content_type: Optional[str]) -> Optional[str]:
        if not raw_text:
            return None

        prompt = f"""你是严格 JSON 修复器。请把下面文本修复为一个合法 JSON 对象。

要求：
1) 只输出 JSON 对象本身
2) 不要 markdown，不要解释
3) 保留原字段语义，缺失字段可补合理默认值

原始文本：
{raw_text[:5000]}
"""
        response_format = {"type": "json_object"} if self.provider != "anthropic" else None
        try:
            return await self._call_ai(
                prompt,
                max_tokens=1600,
                system_prompt="你是严格 JSON 语法修复器，输出必须是可被 json.loads 解析的 JSON 对象。",
                temperature=0.0,
                response_format=response_format,
                content_type=route_content_type,
            )
        except Exception:
            return None

    def _build_stream_script_prompt(
        self,
        news_items: List[Dict],
        editorial_brief: Dict[str, Any],
        duration: int = 30,
        style: str = "专业",
        persona: Optional[Dict] = None,
    ) -> str:
        news_details = "\n".join([
            f"{i+1}. 【{n.get('category', '财经')}】{n['title']}\n   来源：{n['source']}"
            for i, n in enumerate(news_items)
        ])
        categories = list(dict.fromkeys(n.get("category", "财经") for n in news_items))
        category_str = "、".join(categories[:4])
        time_per_news = max(3, duration // max(len(news_items), 1))
        thesis = editorial_brief.get("thesis", editorial_brief.get("lead_angle", ""))
        opening_hook = editorial_brief.get("opening_hook", "")
        contrarian_take = editorial_brief.get("contrarian_take", "")
        core_conflicts = "\n".join(f"- {item}" for item in editorial_brief.get("core_conflicts", []))
        causal_chain = "\n".join(f"- {item}" for item in editorial_brief.get("causal_chain", []))
        market_pulse = "\n".join(f"- {item}" for item in editorial_brief.get("market_pulse", []))
        winners = "\n".join(f"- {item}" for item in editorial_brief.get("winners_losers", {}).get("winners", []))
        losers = "\n".join(f"- {item}" for item in editorial_brief.get("winners_losers", {}).get("losers", []))
        verification_signals = "\n".join(f"- {item}" for item in editorial_brief.get("verification_signals", []))
        takeaways = "\n".join(f"- {item}" for item in editorial_brief.get("audience_takeaways", []))

        emotion_arc_section = ""
        if duration >= 60:
            emotion_arc_section = f"""
【情绪弧线设计（{duration}分钟版）】
- 0-5分钟：用最强冲突或反直觉结论抓住观众
- 5-15分钟：给出第一个高价值“原来如此”时刻
- 15-{duration//2}分钟：铺设背景、逻辑和资金/产业链解释
- {duration//2}-{duration}分钟：抛出核心判断、风险点和明日验证指标
"""

        return f"""【内容任务】
请生成一份真正可上播的财经直播稿，而不是新闻摘要。
{self._persona_section(persona)}
{QUALITY_CONTROL_RULES.strip()}

【今晚最重要的判断】
{thesis or editorial_brief.get('lead_angle', '')}

【推荐开场钩子】
{opening_hook or '先用一个具体数字、反常识事实或市场误判切入'}

【必须讲出的反直觉判断】
{contrarian_take or '市场最容易误判的，不是新闻本身，而是它背后的定价逻辑。'}
【编辑部主线】
{editorial_brief.get('lead_angle', '')}

【你必须抓住的矛盾/预期差】
{core_conflicts or '- 没有额外补充'}

【你必须融入的市场脉搏判断】
{market_pulse or '- 没有额外补充'}

【你必须讲清的因果链】
{causal_chain or '- 从事件变化讲到市场如何重新定价'}

【谁受益、谁受损】
受益方：
{winners or '- 受益方待结合内容判断'}
受损方：
{losers or '- 受损方待结合内容判断'}

【明天最该盯的验证指标】
{verification_signals or '- 暂无额外补充'}

【观众最终要带走的结论】
{takeaways or '- 没有额外补充'}

【今日要闻】共{len(news_items)}条，涉及领域：{category_str}
{news_details}

【直播约束】
- 总时长：约{duration}分钟
- 每条新闻：约{time_per_news}分钟
- 风格：{style}
- 写作气质：{self._style_profile(style)}
{emotion_arc_section}

【结构要求】
1. 开场不要寒暄，前3段内必须完成：给判断、给原因、给观众为什么要听。
2. 每条新闻必须回答：发生了什么、为什么是现在、影响谁、接下来怎么判断。
3. 至少加入 1 个反直觉观点、1 个历史类比、2-3 个生活化比喻。
4. 互动提问必须有信息含量，不是无意义暖场。
5. 结尾必须总结主线、风险点、明日关注指标，并给出一句能被记住的收束判断。
6. 可以加入[互动]、[留人钩子]、[揭晓悬念]、[情绪提示：...]等标记，但不要滥用。

【硬性指标】
- 总字数控制在1500-2200字
- 至少出现4个具体数字
- 禁止空话：建议关注、可能有影响、需要观察、存在不确定性、仅供参考

请直接输出完整直播稿，不要写解释：
"""

    async def generate_stream_script(self, news_items: List[Dict], duration: int = 30, style: str = "专业", persona: Optional[Dict] = None) -> str:
        editorial_brief = await self._prepare_editorial_brief(
            news_items,
            goal="输出一份适合直播口播、判断明确、能帮观众抓主线的财经直播稿",
            style=style,
            route_content_type="stream_script",
        )
        prompt = self._build_stream_script_prompt(news_items, editorial_brief, duration=duration, style=style, persona=persona)
        try:
            response = await self._call_ai(
                prompt,
                max_tokens=3200,
                system_prompt=STREAM_SCRIPT_SYSTEM_PROMPT,
                temperature=0.65,
                content_type="stream_script",
            )
            return self._format_stream_script(response, news_items, duration)
        except Exception as e:
            logger.error("Generation error: %s | %s", e, self._error_context())
            return self._generate_fallback_script(news_items)

    async def stream_stream_script(self, news_items: List[Dict], duration: int = 30, style: str = "专业", persona: Optional[Dict] = None) -> AsyncIterator[str]:
        editorial_brief = await self._prepare_editorial_brief(
            news_items,
            goal="输出一份适合直播口播、判断明确、能帮观众抓主线的财经直播稿",
            style=style,
            route_content_type="stream_script",
        )
        prompt = self._build_stream_script_prompt(news_items, editorial_brief, duration=duration, style=style, persona=persona)
        try:
            async for chunk in self._stream_ai(
                prompt,
                max_tokens=3200,
                system_prompt=STREAM_SCRIPT_SYSTEM_PROMPT,
                content_type="stream_script",
            ):
                if chunk:
                    yield chunk
        except Exception as e:
            logger.error("Streaming generation error: %s | %s", e, self._error_context())
            fallback = self._generate_fallback_script(news_items)
            for chunk in self._chunk_text(fallback, chunk_size=160):
                yield chunk

    async def generate_article(self, news_items: List[Dict], title: str = "", focus_topic: str = "", persona: Optional[Dict] = None) -> Dict:
        editorial_brief = await self._prepare_editorial_brief(
            news_items,
            goal="写一篇能发在公众号上的高质量财经文章，不止要讲发生了什么，更要讲为什么重要",
            style="洞察" if title else "解读型",
            route_content_type="article",
        )
        preferred_title = f"\n【标题方向偏好】\n{title}\n" if title else ""
        focus_hint = f"\n【重点聚焦】\n{focus_topic}\n" if focus_topic else ""
        thesis = editorial_brief.get("thesis", editorial_brief.get("lead_angle", ""))
        contrarian_take = editorial_brief.get("contrarian_take", "")
        article_angles = "\n".join(f"- {item}" for item in editorial_brief.get("article_angles", []))
        winners = "\n".join(f"- {item}" for item in editorial_brief.get("winners_losers", {}).get("winners", []))
        losers = "\n".join(f"- {item}" for item in editorial_brief.get("winners_losers", {}).get("losers", []))
        verification_signals = "\n".join(f"- {item}" for item in editorial_brief.get("verification_signals", []))
        prompt = f"""【内容任务】
请写一篇让人读完就想转发的财经公众号长文。
{preferred_title}{focus_hint}
{self._persona_section(persona)}
{QUALITY_CONTROL_RULES.strip()}

【一句最硬的结论】
{thesis}
【总论点】
{editorial_brief.get('lead_angle', '')}

【最值得写进文章的反直觉判断】
{contrarian_take or '大多数人盯着表面热闹，但真正影响市场的是更深层的主线变化。'}

【必须回答的关键问题】
{chr(10).join(f'- {item}' for item in editorial_brief.get('core_conflicts', []))}

【推荐展开的小节角度】
{article_angles or '- 暂无额外补充'}

【谁会先受益、谁会先承压】
受益方：
{winners or '- 暂无额外补充'}
受损方：
{losers or '- 暂无额外补充'}

【后续最该盯的验证指标】
{verification_signals or '- 暂无额外补充'}

【普通读者最该带走的结论】
{chr(10).join(f'- {item}' for item in editorial_brief.get('audience_takeaways', []))}

【新闻池】
{self._news_snapshot(news_items)}

【写作要求】
1. 必须提供3个标题备选：数字悬念型 / 反直觉型 / 读者利益型。
2. 开篇不能从“今天市场”或“据报道”开始，要用具体场景、反常识事实或数字切入，并在前两段内亮出主判断。
3. 正文按“结论 -> 背景 -> 深层原因 -> 影响 -> 行动”来展开。
4. 每一节都要写清：是什么、为什么、深层原因、明确判断、验证指标。
5. 结尾不能是“点赞关注”，而要给可执行动作、核心判断或下一步验证方法。
6. 不能编造数据、机构观点或内幕。
7. 可以有锋利判断，但必须能自圆其说；不要只写成新闻扩写或摘要拼盘。

【篇幅】
- 总字数控制在1400-2200字
- 至少4个具体数字
- 至少1处历史类比
- 至少1个反主流判断

【输出格式】
===标题选项===
标题A：[内容]
标题B：[内容]
标题C：[内容]

===正文===
（从这里开始输出完整正文，使用Markdown格式）

*本文内容仅供参考，不构成投资建议。投资有风险，决策需谨慎。*
===
"""
        try:
            response = await self._call_ai(
                prompt,
                max_tokens=3200,
                system_prompt=ARTICLE_SYSTEM_PROMPT,
                temperature=0.6,
                content_type="article",
            )
            return self._format_article(response)
        except Exception as e:
            logger.error("Article generation error: %s | %s", e, self._error_context())
            return self._generate_fallback_article(news_items)

    async def generate_flash_report(self, news_items: List[Dict], focus_topic: str = "") -> str:
        editorial_brief = await self._prepare_editorial_brief(
            news_items,
            goal="输出一份可直接发布的财经快报速评，帮助读者30秒抓住主线和判断",
            style="洞察",
        )
        thesis = editorial_brief.get("thesis", editorial_brief.get("lead_angle", ""))
        contrarian_take = editorial_brief.get("contrarian_take", "")
        verification_signals = "\n".join(f"- {item}" for item in editorial_brief.get("verification_signals", []))
        takeaways = "\n".join(f"- {item}" for item in editorial_brief.get("audience_takeaways", []))
        news_details = "\n".join([
            f"{i+1}. 【{n.get('category', '财经')}】{n['title']}（{n['source']}）"
            for i, n in enumerate(news_items[:5])
        ])

        prompt = f"""【快报速评任务】
请输出一篇适合社媒/群发/早会速读的财经快报速评。

{QUALITY_CONTROL_RULES.strip()}

【一句最硬的结论】
{thesis}

【反直觉判断】
{contrarian_take or '真正重要的不是表面热度，而是主线是否开始被重新定价。'}

【素材新闻】
{news_details}

【普通读者最该带走的结论】
{takeaways or '- 暂无额外补充'}

【接下来最该盯的验证指标】
{verification_signals or '- 暂无额外补充'}

【输出结构】
1. 开头：1句话直接给最重要判断
2. 速评主体：按3-5个短段落写，每段都要说“这意味着什么”
3. 结尾：一句短金句 + 1个下一步验证点

【硬性要求】
- 总字数控制在350-500字
- 至少出现2个具体数字或时间点
- 不能写成摘要拼盘
- 不能使用“建议关注 / 需要观察 / 影响深远”之类空话
- 结尾金句不超过18字，能被截图传播

请直接输出快报速评正文，不要解释：
"""
        try:
            response = await self._call_ai(
                prompt,
                max_tokens=1200,
                system_prompt=FLASH_REPORT_SYSTEM_PROMPT,
            )
            return self._clean_generated_text(response)
        except Exception as e:
            logger.error("Flash report generation error: %s | %s", e, self._error_context())
            return self._generate_fallback_flash_report(news_items)

    async def _analyze_news_outline(self, news_items: List[Dict], focus_topic: str, date_str: str) -> str:
        news_details = "\n".join([f"• {n['title']}（{n['source']}）" for n in news_items[:6]])
        analysis_prompt = f"""请对以下财经新闻做内部分析（不写文章，只做分析框架）：

核心事件：{focus_topic}
相关新闻：
{news_details}
分析日期：{date_str}

请按以下结构输出简洁分析框架（不超过600字）：
1. 反共识洞察：主流判断是什么？被低估的一面是什么？
2. 利益流向图：谁受益，谁受损，有哪些意外受益者？
3. 历史镜像：最相似的历史案例 + 关键不同点
4. 二阶影响：列出3个大多数分析忽视的间接影响
5. 核心金句：一句话总结本次事件的本质（≤20字）
6. 6个月预测：一个具体可量化可验证的预测
"""
        try:
            return await self._call_ai(
                analysis_prompt,
                max_tokens=1200,
                system_prompt=DEEP_DIVE_SYSTEM_PROMPT,
            )
        except Exception as e:
            logger.warning("Outline analysis failed, proceeding without: %s", e)
            return ""

    async def generate_deep_dive(self, news_items: List[Dict], focus_topic: str = "", persona: Optional[Dict] = None) -> str:
        if not focus_topic:
            focus_topic = news_items[0]["title"]

        editorial_brief = await self._prepare_editorial_brief(
            news_items,
            goal="写一篇有研究框架、有投资视角、能帮助读者建立判断的深度长文",
            style="洞察",
        )
        date_str = datetime.now().strftime("%Y年%m月%d日")
        internal_outline = await self._analyze_news_outline(news_items, focus_topic, date_str)
        outline_context = f"\n【预分析框架】\n{internal_outline}\n" if internal_outline else ""
        prompt = f"""【深度研究任务】
你要写的不是新闻汇总，而是一篇能被人截图传播的原创深度分析。
{self._persona_section(persona)}
【总论点】
{editorial_brief.get('lead_angle', '')}
{outline_context}
【核心事件】
{focus_topic}

【新闻池】
{self._news_snapshot(news_items[:10])}

【必须拆清楚的问题】
{chr(10).join(f'- {item}' for item in editorial_brief.get('core_conflicts', []))}

【写作结构】
## 一、破题——打破认知
## 二、溯源——这件事为什么走到今天
## 三、深度拆解——从宏观/产业/资金/情绪看问题
## 四、历史镜鉴——上次类似的事怎么演化
## 五、蝴蝶效应——二阶和三阶影响
## 六、情景推演——乐观/中性/悲观三种路径与概率
## 七、先行指标——接下来该盯什么验证
## 八、核心结论——6个月内的可验证预测

【强制要求】
- 总字数不少于3000字
- 至少10个具体数字
- 三种情景必须给出具体概率，且和为100%
- 至少1个反主流判断
- 每章节之间要有桥接句，不要硬切换
- 结尾必须给出6个月内可验证预测
- 禁止空话：建议关注、可能有影响、需要观察、存在不确定性、仅供参考

请直接输出完整深度长文，使用Markdown格式：
"""
        try:
            return await self._call_ai(prompt, max_tokens=7000, system_prompt=DEEP_DIVE_SYSTEM_PROMPT)
        except Exception as e:
            logger.error("Deep dive generation error: %s | %s", e, self._error_context())
            return self._generate_fallback_deep_dive(news_items)

    async def generate_ppt_script(self, news_items: List[Dict], focus_topic: str = "", persona: Optional[Dict] = None) -> str:
        if not focus_topic:
            focus_topic = news_items[0]["title"]
        editorial_brief = await self._prepare_editorial_brief(
            news_items,
            goal="输出一份适合晨会、直播前准备或路演汇报的财经 PPT 脚本",
            style="洞察",
        )
        news_details = "\n".join([
            f"{i+1}. 【{n.get('category','财经')}】{n['title']}（{n['source']}）"
            for i, n in enumerate(news_items[:6])
        ])
        prompt = f"""【PPT创作任务】
你要生成一套可直接用于财经讲解、路演、直播准备的完整 PPT 脚本。
{self._persona_section(persona)}
【核心主题】
{focus_topic}

【总论点】
{editorial_brief.get('lead_angle', '')}

【必须传递给观众的结论】
{chr(10).join(f'- {item}' for item in editorial_brief.get('audience_takeaways', []))}

【素材新闻】
{news_details}

【输出要求】
请按页输出，每页都必须包含以下字段：
- 【第X张 · 页面类型】
- 标题：
- 💡 核心论点：
- 📌 屏幕要点：3-4条
- 🎤 讲者逐字稿：120-220字
- 📊 可视化建议：

【页面建议结构】
1. 封面
2. 结论前置：今天最重要的3个判断
3. 为什么这件事比你想的更重要
4. 打破认知：主流判断哪里错了
5-N. 逐条新闻深析
N+1. 把几条新闻串起来的共同逻辑
N+2. 情景推演
N+3. 行动指南
N+4. 记忆点与先行指标

【硬性要求】
- 每页核心论点必须是观点句，不是描述句
- 屏幕要点保持极简，每条不超过18字
- 口播要像真人在讲，不是官方发言稿
- 至少4页出现具体数字
- 至少1页是反直觉判断
- 禁止空话：建议关注、可能有影响、需要观察、存在不确定性、仅供参考

请直接输出完整PPT脚本：
"""
        try:
            return await self._call_ai(prompt, max_tokens=6500, system_prompt=PPT_SYSTEM_PROMPT)
        except Exception as e:
            logger.error("PPT script generation error: %s | %s", e, self._error_context())
            return self._generate_fallback_ppt_script(news_items, focus_topic)

    # ─────────────────────────────────────────────────────────────────
    # 核心壁垒 1：财经风控合规 Agent
    # ─────────────────────────────────────────────────────────────────

    async def compliance_review(self, content: str) -> Dict[str, Any]:
        """对已生成内容做财经合规审核，返回风险报告和改写版本。"""
        if not content or not content.strip():
            return {
                "is_compliant": True,
                "risk_level": "低",
                "issues": [],
                "revised_content": content,
                "summary": "内容为空，无需审核。",
            }
        prompt = f"""请对以下财经内容进行全面合规审核，识别所有违规表达并输出改写版本。

【待审核内容】
{content[:6000]}

请严格按照系统要求的 JSON 格式输出，不要添加任何额外说明。
"""
        try:
            raw = await self._call_ai(prompt, max_tokens=4000, system_prompt=COMPLIANCE_SYSTEM_PROMPT)
            result = self._extract_json(raw)
            # Ensure required fields exist
            result.setdefault("is_compliant", True)
            result.setdefault("risk_level", "低")
            result.setdefault("issues", [])
            result.setdefault("revised_content", content)
            result.setdefault("summary", "合规审核完成。")
            return result
        except Exception as e:
            logger.error("Compliance review error: %s | %s", e, self._error_context())
            return {
                "is_compliant": True,
                "risk_level": "低",
                "issues": [],
                "revised_content": content,
                "summary": "合规审核服务暂时不可用，请人工复核。",
            }

    # ─────────────────────────────────────────────────────────────────
    # 核心壁垒 3：一键内容矩阵（朋友圈预热 + 直播稿 + 复盘文章 + PPT脚本）
    # ─────────────────────────────────────────────────────────────────

    async def generate_moments_copy(
        self, news_items: List[Dict], focus_topic: str = "", live_time: str = "", persona: Optional[Dict] = None
    ) -> str:
        """生成朋友圈/微博预热诱饵文案（50字以内）。"""
        topic = focus_topic or (news_items[0]["title"] if news_items else "今日财经要闻")
        time_hint = f"（晚{live_time}点直播）" if live_time else "（今晚直播）"
        prompt = f"""请根据以下财经主题，写一条50字以内的朋友圈预热文案，用于吸引粉丝晚上来看直播。
{self._persona_section(persona)}
【今日核心话题】{topic}
【直播时间提示】{time_hint}
【相关新闻】{self._news_snapshot(news_items[:3])}

直接输出文案正文，不要标题，不要解释，不超过50字：
"""
        try:
            return (await self._call_ai(prompt, max_tokens=200, system_prompt=MOMENTS_COPY_SYSTEM_PROMPT)).strip()
        except Exception as e:
            logger.error("Moments copy generation error: %s | %s", e, self._error_context())
            return f"今天有一件事，可能会让你的判断彻底改变。{time_hint}，直播间见。"

    async def generate_content_matrix(
        self,
        news_items: List[Dict],
        focus_topic: str = "",
        duration: int = 30,
        style: str = "专业",
        live_time: str = "",
        persona: Optional[Dict] = None,
    ) -> Dict[str, Any]:
        """一键生成完整内容矩阵：朋友圈预热 + 直播稿 + 复盘文章 + PPT脚本。"""
        moments_task = asyncio.create_task(
            self.generate_moments_copy(news_items, focus_topic=focus_topic, live_time=live_time, persona=persona)
        )
        script_task = asyncio.create_task(
            self.generate_stream_script(news_items, duration=duration, style=style, persona=persona)
        )
        article_task = asyncio.create_task(
            self.generate_article(news_items, focus_topic=focus_topic, persona=persona)
        )
        ppt_task = asyncio.create_task(
            self.generate_ppt_script(news_items, focus_topic=focus_topic, persona=persona)
        )

        moments, stream_script, article, ppt_script = await asyncio.gather(
            moments_task, script_task, article_task, ppt_task
        )

        return {
            "moments_copy": moments,
            "stream_script": stream_script,
            "article": article,
            "ppt_script": ppt_script,
            "generated_at": datetime.now().isoformat(),
        }

    async def generate_ppt(self, news_items: List[Dict], title: str = "", style: str = "专业", focus_topic: str = "") -> bytes:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        ppt_script = await self.generate_ppt_script(news_items, focus_topic=focus_topic or title)
        slides = self._parse_ppt_script(ppt_script)
        if not slides:
            slides = self._fallback_slide_data(news_items, title or focus_topic)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6]

        blue = RGBColor(0x1A, 0x56, 0xDB)
        dark = RGBColor(0x1F, 0x2A, 0x3C)
        white = RGBColor(0xFF, 0xFF, 0xFF)
        gray = RGBColor(0x6B, 0x72, 0x80)

        def add_bg(slide, color):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
            tx_box = slide.shapes.add_textbox(left, top, width, height)
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            return tx_box

        cover_title = title or focus_topic or f"财经简报 · {datetime.now().strftime('%Y年%m月%d日')}"
        cover = prs.slides.add_slide(layout)
        add_bg(cover, dark)
        add_textbox(cover, cover_title, Inches(0.9), Inches(1.7), Inches(11.4), Inches(1.5), font_size=32, bold=True, color=white, align=PP_ALIGN.CENTER)
        add_textbox(cover, slides[0].get('core', '今天最重要的判断，先讲结论再讲论证'), Inches(1.1), Inches(3.4), Inches(11.1), Inches(1), font_size=18, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)

        for slide_data in slides:
            slide = prs.slides.add_slide(layout)
            add_bg(slide, white)
            add_textbox(slide, slide_data['title'], Inches(0.7), Inches(0.5), Inches(11.8), Inches(0.8), font_size=24, bold=True, color=dark)
            add_textbox(slide, f"💡 核心论点：{slide_data['core']}", Inches(0.8), Inches(1.4), Inches(11.4), Inches(0.8), font_size=15, bold=True, color=blue)
            add_textbox(slide, "📌 屏幕要点", Inches(0.8), Inches(2.2), Inches(3), Inches(0.4), font_size=13, bold=True, color=blue)
            top = 2.7
            for bullet in slide_data['bullets'][:4]:
                add_textbox(slide, f"• {bullet}", Inches(1.0), Inches(top), Inches(5.4), Inches(0.45), font_size=14, color=dark)
                top += 0.45
            add_textbox(slide, "🎤 讲者备注", Inches(6.4), Inches(2.2), Inches(3), Inches(0.4), font_size=13, bold=True, color=blue)
            add_textbox(slide, slide_data['speaker_notes'], Inches(6.5), Inches(2.7), Inches(5.7), Inches(2.7), font_size=13, color=dark)
            add_textbox(slide, f"📊 可视化建议：{slide_data['visual']}", Inches(0.8), Inches(5.9), Inches(11.2), Inches(0.6), font_size=12, color=gray)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _parse_ppt_script(self, script: str) -> List[Dict[str, Any]]:
        sections = re.split(r"(?=【第\d+张)", script or "")
        slides: List[Dict[str, Any]] = []
        for section in sections:
            section = section.strip()
            if not section.startswith("【第"):
                continue
            title = ""
            core = ""
            bullets: List[str] = []
            speaker_notes = ""
            visual = ""

            title_match = re.search(r"标题[:：]\s*(.+)", section)
            if title_match:
                title = title_match.group(1).strip()
            else:
                header_match = re.search(r"【第\d+张[^】]*】\s*(.+)", section)
                title = header_match.group(1).strip() if header_match else "未命名页"

            core_match = re.search(r"[💡🎯]\s*(?:核心论点|核心信息)[:：]\s*(.+)", section)
            if core_match:
                core = core_match.group(1).strip()

            bullet_match = re.search(r"📌\s*(?:屏幕上的文字要点|屏幕要点|要点)[:：](.*?)(?:🎤|📊|$)", section, re.DOTALL)
            if bullet_match:
                for line in bullet_match.group(1).splitlines():
                    line = re.sub(r"^[\s•·\-]+", "", line).strip()
                    if line:
                        bullets.append(line)

            speaker_match = re.search(r"🎤\s*(?:讲者逐字稿|演讲者口播)[^\n]*\n(.*?)(?:📊|$)", section, re.DOTALL)
            if speaker_match:
                speaker_notes = "\n".join(line.strip() for line in speaker_match.group(1).splitlines() if line.strip())

            visual_match = re.search(r"📊\s*(?:可视化建议|视觉设计建议)[:：]\s*(.+)", section)
            if visual_match:
                visual = visual_match.group(1).strip()

            if title or core or bullets or speaker_notes:
                slides.append({
                    "title": title or "未命名页",
                    "core": core or "这一页的核心是帮助观众抓住真正的主线",
                    "bullets": bullets or ["补充关键事实", "解释深层逻辑", "提示后续验证指标"],
                    "speaker_notes": speaker_notes or "这一页要讲清楚：为什么这件事重要、为什么是现在、接下来要看什么。",
                    "visual": visual or "建议使用对比结构、关键数字放大或时间线布局。",
                })
        return slides

    def _fallback_slide_data(self, news_items: List[Dict], focus_topic: str = "") -> List[Dict[str, Any]]:
        topic = focus_topic or news_items[0]["title"]
        slides = [{
            "title": "今天最重要的3个判断",
            "core": "先讲结论，再讲论证，帮助观众快速抓住主线。",
            "bullets": [
                "市场真正定价的是主线变化",
                "短期情绪不等于中期趋势",
                "接下来要看验证指标",
            ],
            "speaker_notes": "先把答案告诉观众，再逐步解释为什么，避免观众在前3分钟流失。",
            "visual": "建议使用三栏对比布局，突出结论1/2/3。",
        }]
        for news in news_items[:6]:
            slides.append({
                "title": news["title"][:26],
                "core": f"这条新闻真正重要的地方，在于它会改变{news.get('category', '财经')}方向的预期。",
                "bullets": [
                    f"来源：{news['source']}",
                    f"分类：{news.get('category', '财经')}",
                    "看清市场真正关心什么",
                ],
                "speaker_notes": f"这条新闻表面上是{news.get('category', '财经')}事件，真正要讲的是它如何影响预期、情绪和后续验证指标。",
                "visual": "建议采用标题+三条 bullet + 底部备注的简洁结构。",
            })
        slides.append({
            "title": topic[:26],
            "core": "如果只能记住一件事，那就是主线比碎片更重要。",
            "bullets": ["看主线", "看验证", "看风险收益比"],
            "speaker_notes": "收尾时把整场逻辑合成一句能被记住的话，并提醒观众明天开始该看什么。",
            "visual": "建议极简收尾页，核心句居中，验证指标置于下方。",
        })
        return slides

    async def _call_ai(
        self,
        prompt: str,
        max_tokens: int = 4000,
        system_prompt: Optional[str] = None,
        temperature: float = 0.75,
        response_format: Optional[Dict[str, str]] = None,
        content_type: Optional[str] = None,
    ) -> str:
        sys_prompt = system_prompt or MASTER_SYSTEM_PROMPT
        last_error: Optional[Exception] = None

        if self._should_use_quality_router(content_type):
            return await self._call_openrouter(
                prompt,
                max_tokens=max_tokens,
                system_prompt=sys_prompt,
                temperature=temperature,
                content_type=content_type,
                response_format=response_format,
            )

        if not self.api_key:
            raise RuntimeError(f"缺少 {self.provider} 的 API Key")

        for attempt in range(2):
            try:
                if self.provider == "anthropic":
                    message = await self.client.messages.create(
                        model=self.model,
                        max_tokens=max_tokens,
                        temperature=temperature,
                        system=sys_prompt,
                        messages=[{"role": "user", "content": prompt}],
                    )
                    usage = self._extract_usage_tokens(getattr(message, "usage", None))
                    self._record_usage_event(
                        provider=self.provider,
                        model=self.model,
                        content_type=content_type,
                        prompt_tokens=usage["prompt_tokens"],
                        completion_tokens=usage["completion_tokens"],
                        route="base",
                    )
                    return message.content[0].text

                request_kwargs: Dict[str, Any] = {
                    "model": self.model,
                    "messages": [
                        {"role": "system", "content": sys_prompt},
                        {"role": "user", "content": prompt},
                    ],
                    "max_tokens": max_tokens,
                    "temperature": temperature,
                }
                if response_format:
                    request_kwargs["response_format"] = response_format

                response = await self.client.chat.completions.create(**request_kwargs)
                usage = self._extract_usage_tokens(getattr(response, "usage", None))
                self._record_usage_event(
                    provider=self.provider,
                    model=self.model,
                    content_type=content_type,
                    prompt_tokens=usage["prompt_tokens"],
                    completion_tokens=usage["completion_tokens"],
                    route="base",
                )
                content = response.choices[0].message.content
                return content if content else ""
            except Exception as exc:
                last_error = exc
                if attempt == 0:
                    await asyncio.sleep(1.2)
                    continue
                raise last_error

        raise last_error if last_error else RuntimeError("unknown ai call error")

    async def _call_openrouter(
        self,
        prompt: str,
        max_tokens: int,
        system_prompt: str,
        temperature: float,
        content_type: str,
        response_format: Optional[Dict[str, str]] = None,
    ) -> str:
        models = self._quality_route_models(content_type)
        if not models or self.openrouter_client is None:
            raise RuntimeError("OpenRouter 质量路由未正确配置")

        request_kwargs: Dict[str, Any] = {
            "model": models[0],
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt},
            ],
            "max_tokens": max_tokens,
            "temperature": temperature,
            "extra_body": {
                "models": models,
                "route": "fallback",
            },
        }
        if response_format:
            request_kwargs["response_format"] = response_format

        response = await self.openrouter_client.chat.completions.create(**request_kwargs)
        selected_model = getattr(response, "model", models[0])
        usage = self._extract_usage_tokens(getattr(response, "usage", None))
        self._record_usage_event(
            provider="openrouter",
            model=selected_model,
            content_type=content_type,
            prompt_tokens=usage["prompt_tokens"],
            completion_tokens=usage["completion_tokens"],
            route="quality",
        )
        logger.info(
            "OpenRouter routed content_type=%s selected_model=%s candidate_models=%s prompt_tokens=%s completion_tokens=%s",
            content_type,
            selected_model,
            ",".join(models),
            usage["prompt_tokens"],
            usage["completion_tokens"],
        )
        content = response.choices[0].message.content
        return content if content else ""

    async def _stream_ai(
        self,
        prompt: str,
        max_tokens: int = 2000,
        system_prompt: Optional[str] = None,
        content_type: Optional[str] = None,
    ) -> AsyncIterator[str]:
        sys_prompt = system_prompt or MASTER_SYSTEM_PROMPT

        if self._should_use_quality_router(content_type):
            async for chunk in self._stream_openrouter(
                prompt,
                max_tokens=max_tokens,
                system_prompt=sys_prompt,
                content_type=content_type,
            ):
                yield chunk
            return

        if not self.api_key:
            raise RuntimeError(f"缺少 {self.provider} 的 API Key")

        if self.provider == "anthropic":
            full_text = await self._call_ai(prompt, max_tokens=max_tokens, system_prompt=sys_prompt)
            for chunk in self._chunk_text(full_text):
                yield chunk
            return

        stream = await self.client.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": prompt},
            ],
            max_tokens=max_tokens,
            temperature=0.75,
            stream=True,
        )
        async for part in stream:
            delta = ""
            if part.choices:
                delta = part.choices[0].delta.content or ""
            if delta:
                yield delta

    async def _stream_openrouter(
        self,
        prompt: str,
        max_tokens: int,
        system_prompt: str,
        content_type: str,
    ) -> AsyncIterator[str]:
        models = self._quality_route_models(content_type)
        if not models or self.openrouter_client is None:
            raise RuntimeError("OpenRouter 质量路由未正确配置")

        stream = await self.openrouter_client.chat.completions.create(
            model=models[0],
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt},
            ],
            max_tokens=max_tokens,
            temperature=0.65,
            stream=True,
            stream_options={"include_usage": True},
            extra_body={
                "models": models,
                "route": "fallback",
            },
        )
        selected_model = models[0]
        prompt_tokens = 0
        completion_tokens = 0
        async for part in stream:
            part_model = getattr(part, "model", None)
            if part_model:
                selected_model = part_model
            usage = self._extract_usage_tokens(getattr(part, "usage", None))
            if usage["prompt_tokens"] or usage["completion_tokens"]:
                prompt_tokens = usage["prompt_tokens"]
                completion_tokens = usage["completion_tokens"]
            delta = ""
            if part.choices:
                delta = part.choices[0].delta.content or ""
            if delta:
                yield delta
        self._record_usage_event(
            provider="openrouter",
            model=selected_model,
            content_type=content_type,
            prompt_tokens=prompt_tokens,
            completion_tokens=completion_tokens,
            route="quality_stream",
        )
        logger.info(
            "OpenRouter streamed content_type=%s selected_model=%s candidate_models=%s prompt_tokens=%s completion_tokens=%s",
            content_type,
            selected_model,
            ",".join(models),
            prompt_tokens,
            completion_tokens,
        )

    def _chunk_text(self, text: str, chunk_size: int = 120) -> List[str]:
        return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

    def _clean_generated_text(self, content: str) -> str:
        cleaned = (content or "").strip()
        cleaned = re.sub(r"^\s*(好的|下面是|以下是)[：:\s]*", "", cleaned)
        replacements = {
            "建议关注": "重点看",
            "需要观察": "下一步验证",
            "可能有影响": "会直接影响",
            "影响深远": "会重塑后续定价",
            "存在不确定性": "仍待验证",
        }
        for src, dst in replacements.items():
            cleaned = cleaned.replace(src, dst)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        return cleaned.strip()

    def _format_stream_script(self, content: str, news_items: List[Dict], duration: int = 30) -> str:
        content = self._clean_generated_text(content)
        date_str = datetime.now().strftime("%Y年%m月%d日")
        weekday = ["星期一", "星期二", "星期三", "星期四", "星期五", "星期六", "星期日"][datetime.now().weekday()]
        return f"""
╔══════════════════════════════════════════════════════════╗
║  📺 财经深度直播稿                                        ║
║  📅 {date_str} {weekday}                                  ║
║  ⏱ 预计时长：{duration}分钟  📊 共{len(news_items)}条新闻  ║
╚══════════════════════════════════════════════════════════╝

{content}

{'─'*60}
📌 本稿由 AI 深度创作生成，请结合实时盘面和最新数据调整
📊 预计直播时长：{len(news_items) * 4}-{len(news_items) * 6} 分钟
{'─'*60}
"""

    def _format_article(self, content: str) -> Dict:
        content = self._clean_generated_text(content)
        titles = []
        if "===标题选项===" in content:
            parts = content.split("===正文===")
            title_section = parts[0].replace("===标题选项===", "").strip()
            raw_titles = [t.strip() for t in title_section.split("\n") if t.strip()]
            for t in raw_titles:
                cleaned = re.sub(r"^标题[A-Ca-c：:]\s*", "", t).strip()
                if cleaned:
                    titles.append(cleaned)
            titles = titles[:3]
            content_body = parts[1] if len(parts) > 1 else content
            content_body = content_body.split("===")[0].strip()
        else:
            title_match = re.search(r'【标题选项】\n(.+?)\n\n', content, re.DOTALL)
            if title_match:
                titles = [t.strip("- 123. ") for t in title_match.group(1).split("\n") if t.strip()]
            content_body = content

        if not titles:
            titles = [f"财经深度解读 {datetime.now().strftime('%m月%d日')}"]

        return {
            "titles": titles,
            "content": content_body.strip(),
            "html": self._markdown_to_html(content_body),
        }

    def _markdown_to_html(self, md: str) -> str:
        html = md
        html = html.replace("===标题选项===", "").replace("===正文===", "").replace("===", "")
        html = re.sub(r"标题[A-Ca-c：:].+?\n", "", html)
        html = re.sub(r'【标题选项】.+?\n\n', '', html, flags=re.DOTALL)
        html = re.sub(r'^# (.+)$', r'<h1 style="font-size:22px;font-weight:bold;margin:20px 0;color:#1a1a1a;">\1</h1>', html, flags=re.MULTILINE)
        html = re.sub(r'^## (.+)$', r'<h2 style="font-size:18px;font-weight:bold;margin:16px 0 8px;color:#1a1a1a;border-left:4px solid #0066cc;padding-left:10px;">\1</h2>', html, flags=re.MULTILINE)
        html = re.sub(r'^### (.+)$', r'<h3 style="font-size:16px;font-weight:bold;margin:12px 0 6px;color:#333;">\1</h3>', html, flags=re.MULTILINE)
        html = re.sub(r'\*\*(.+?)\*\*', r'<strong style="color:#0066cc;">\1</strong>', html)
        html = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html)
        html = re.sub(r'^[•·▶📍📌✅❌🎯💡❓🟢🟡🔴] (.+)$', r'<li style="margin:6px 0;">\1</li>', html, flags=re.MULTILINE)
        html = re.sub(r'^- (.+)$', r'<li style="margin:6px 0;">\1</li>', html, flags=re.MULTILINE)
        html = re.sub(r'(<li.+</li>\n?)+', lambda m: f'<ul style="margin:10px 0;padding-left:20px;list-style:disc;">{m.group(0)}</ul>', html)
        paragraphs = html.split("\n\n")
        formatted = []
        for p in paragraphs:
            p = p.strip()
            if p and not p.startswith("<h") and not p.startswith("<ul") and not p.startswith("<li"):
                p = f'<p style="line-height:1.9;margin:15px 0;color:#333;font-size:15px;">{p}</p>'
            formatted.append(p)
        return "\n".join(formatted)

    def _generate_fallback_script(self, news_items: List[Dict]) -> str:
        date_str = datetime.now().strftime("%Y年%m月%d日")
        script = f"""
【开场判断】
今天最重要的，不是某一条新闻本身，而是市场正在重新定价的方向。今天是{date_str}，下面我们把真正值得看的主线拆开来说。

"""
        for i, news in enumerate(news_items, 1):
            script += f"""
【新闻{i}】{news['title']}

【深度解读】
这条消息来自{news['source']}，属于{news.get('category', '财经')}。表面上它是一条单点新闻，但更值得关注的是它会如何影响预期、情绪和后续验证指标。

[互动] 如果你也在盯这个方向，在评论区扣1

"""
        script += """
【收尾总结】
把今天这些消息合在一起看，真正该盯的是主线，而不是单条标题带来的短期情绪。

【明日关注】
- 资金流向是否持续
- 关键政策或数据是否跟进验证
- 相关板块是否从情绪走向业绩与基本面
"""
        return script

    def _generate_fallback_article(self, news_items: List[Dict]) -> Dict:
        date_str = datetime.now().strftime("%m月%d日")
        titles = [
            f"今天这{len(news_items)}件事，比大多数人想得更重要",
            f"一个被忽视的信号：{news_items[0]['title'][:18]}背后的逻辑",
            f"{date_str}财经深度：看懂这几件事，看懂接下来的市场",
        ]
        content = f"# {titles[0]}\n\n"
        content += f"**{datetime.now().strftime('%Y年%m月%d日')} | 深度解读**\n\n"
        content += f"今天发生了{len(news_items)}件值得认真对待的事。把它们放在一起看，比逐条看热闹更重要。\n\n"
        for news in news_items:
            content += f"## {news['title']}\n\n"
            content += f"**来源**：{news['source']} | **分类**：{news.get('category', '财经')}\n\n"
            content += "这条新闻的真正价值，不只在于表面事实，而在于它如何改变市场预期、资金定价和后续验证路径。\n\n"
        content += "## 核心判断\n\n把今天这些消息放在一起看，市场真正变化的是主线和预期差，而不只是热闹本身。\n\n"
        content += "*本文内容仅供参考，不构成投资建议。*\n"
        return {"titles": titles, "content": content, "html": self._markdown_to_html(content)}

    def _generate_fallback_flash_report(self, news_items: List[Dict]) -> str:
        lead = news_items[0] if news_items else {"title": "今日财经主线出现新变化", "source": "系统"}
        lines = [
            f"今天最重要的判断是：{lead['title'][:28]}，真正值得看的不是新闻热度，而是它背后的主线变化。",
            "",
        ]
        for news in news_items[:4]:
            lines.append(
                f"【{news.get('category', '财经')}】{news['title']}\n"
                f"这条消息来自{news['source']}，表面上是一条事件新闻，真正重要的是它会改变市场对后续节奏和验证指标的判断。"
            )
            lines.append("")
        lines.append("一句话总结：看主线，不要只看热闹。")
        lines.append("下一步重点看：政策跟进、资金连续性和基本面验证。")
        return "\n".join(lines)

    def _generate_fallback_deep_dive(self, news_items: List[Dict]) -> str:
        topic = news_items[0]['title']
        date_str = datetime.now().strftime("%Y年%m月%d日")
        return f"""# 深度分析：{topic}

**{date_str} | 原创深度研究**

## 一、先说结论
市场最容易被高估的，是短期情绪；最容易被低估的，是主线变化的持续性。

## 二、事件到底改变了什么
这件事改变的不只是表面信息，而是市场对后续政策、行业景气度和资金偏好的预期。

## 三、市场为什么会对它敏感
因为它连接的是预期差，而不是静态事实。

## 四、真正的机会在哪里
机会不在标题最热的地方，而在后续更容易被验证的方向。

## 五、最大的风险与误判点
最大的误判，是把一次性情绪催化当成长期趋势。

## 六、接下来怎么跟踪验证
重点看政策、资金流、订单、价格和盈利验证节奏。

## 七、总结
如果6个月后回头看，今天更可能是一个预期重定价的起点，而不是单日噪音。
"""

    def _generate_fallback_ppt_script(self, news_items: List[Dict], focus_topic: str = "") -> str:
        topic = focus_topic or news_items[0]['title']
        return f"""【第1张 · 封面】
标题：{topic[:20]}——它比你想的更重要
💡 核心论点：今天真正要讲的是主线，而不是碎片新闻
📌 屏幕要点：
• 先看结论
• 再看论证
• 最后看行动
🎤 讲者逐字稿：今天我们不一条条念新闻，而是把这些消息重新放回同一张地图里，看看到底哪条线最值得你花时间。
📊 可视化建议：深色背景 + 大标题 + 一句判断。
"""


generator = ContentGenerator()
